"""Normalize data/clean documents into citation-ready JSONL chunks, and
optionally write them into Course/Document/DocumentChunk (DB) so the RAG
retrieval path can actually surface them.

Usage:
    python backend/scripts/ingest_clean_corpus.py --source data/clean --output data/cursus_corpus.jsonl
    python backend/scripts/ingest_clean_corpus.py --source data/clean --output data/cursus_corpus.jsonl --write-db

Legacy .ppt files are converted through LibreOffice headless before extraction.
A manifest (`<output>.manifest.json`, one entry per source file keyed by
relative path) tracks each file's content checksum + resulting document_id +
chunk_count + status, so a re-run skips files whose content hasn't changed
and can report exactly which files failed and why, instead of silently
skipping them.
"""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
import tempfile
from hashlib import sha256
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

SUPPORTED = {".docx", ".pdf", ".pptx", ".ppt", ".xlsx"}
CHUNK_SIZE = 1100
SOURCE_TAG = "clean_corpus"


def subject_code(path: Path) -> str:
    """Course code from a `<N>.<CODE>` directory segment (e.g. `1.CEA201`).
    Scans directory parts only (`path.parts[:-1]`) -- the filename itself
    almost always contains a dot too (its extension, or an embedded one like
    "ASP.NET"), and scanning it produced false-positive "codes" like "pptx"
    or "docx" whenever the part before the LAST dot happened to be alnum."""
    for part in path.parts[:-1]:
        if "." in part and part.split(".", 1)[-1].upper().isalnum():
            return part.split(".", 1)[-1]
    return path.parent.name


def extract(path: Path) -> list[tuple[str, str]]:
    if path.suffix == ".docx":
        doc = DocxDocument(path)
        return [("document", "\n".join(p.text for p in doc.paragraphs if p.text.strip()))]
    if path.suffix == ".pdf":
        return [(f"page {i}", page.extract_text() or "") for i, page in enumerate(PdfReader(path).pages, 1)]
    if path.suffix == ".ppt":
        if not shutil.which("soffice"):
            raise RuntimeError("LibreOffice (soffice) is required for legacy .ppt")
        with tempfile.TemporaryDirectory() as temp:
            subprocess.run(["soffice", "--headless", "--convert-to", "pptx", "--outdir", temp, str(path)], check=True, capture_output=True)
            converted = Path(temp) / f"{path.stem}.pptx"
            return extract(converted)
    if path.suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            sheets = []
            for sheet in workbook.worksheets:
                lines = [
                    "\t".join(str(cell) for cell in row if cell is not None)
                    for row in sheet.iter_rows(values_only=True)
                    if any(cell is not None for cell in row)
                ]
                sheets.append((f"sheet {sheet.title}", "\n".join(lines)))
            return sheets
        finally:
            workbook.close()
    deck = Presentation(path)
    return [(f"slide {i}", "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())) for i, slide in enumerate(deck.slides, 1)]


def chunk_text(text: str, size: int = CHUNK_SIZE) -> list[str]:
    text = text.strip()
    return [piece for i in range(0, len(text), size) if (piece := text[i : i + size].strip())]


def _file_checksum(path: Path) -> str:
    return sha256(path.read_bytes()).hexdigest()


def _load_manifest(path: Path) -> dict:
    if not path.is_file():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}


def _slugify_document_id(rel_path: str) -> str:
    return "doc_clean_" + sha256(rel_path.encode()).hexdigest()[:20]


def write_db(*, rel_path: str, checksum: str, subject: str, title: str, sections: list[tuple[str, str]]) -> tuple[str, int]:
    """Upsert one Course + one Document (one per source file) + its
    DocumentChunks. Returns (document_id, chunk_count). Idempotent: deletes
    and re-inserts this document's chunks every call (cheap — a document's
    chunk count tops out in the low hundreds), matching
    `ingest_subject_data.py::write_db`'s convention."""
    from src.db import models
    from src.db.connection import SessionLocal

    document_id = _slugify_document_id(rel_path)
    db = SessionLocal()
    try:
        course = db.query(models.Course).filter_by(code=subject).first()
        if course is None:
            course = models.Course(id=subject, code=subject, name=subject, description=subject)
            db.add(course)
            db.flush()

        document = db.query(models.Document).filter_by(id=document_id).first()
        metadata_info = {"source": SOURCE_TAG, "course_code": subject, "original_path": rel_path, "checksum": checksum}
        if document is None:
            document = models.Document(
                id=document_id,
                course_id=course.id,
                title=title,
                file_path=rel_path,
                doc_type="LECTURE",
                version="1.0",
                metadata_info=metadata_info,
            )
            db.add(document)
            db.flush()
        else:
            document.title = title
            document.file_path = rel_path
            document.metadata_info = metadata_info

        db.query(models.DocumentChunk).filter_by(document_id=document.id).delete(synchronize_session=False)

        chunk_count = 0
        for location, text in sections:
            for piece in chunk_text(text):
                db.add(
                    models.DocumentChunk(
                        id=f"chunk_{document.id}_{chunk_count}",
                        document_id=document.id,
                        chunk_index=chunk_count,
                        text=piece,
                        token_count=max(1, len(piece.split())),
                        metadata_info={
                            "course_code": subject,
                            "doc_type": "LECTURE",
                            "doc_title": title,
                            "section": location,
                            "source_label": f"{title} · {location}",
                            "source": SOURCE_TAG,
                        },
                    )
                )
                chunk_count += 1
        db.commit()
        return document_id, chunk_count
    finally:
        db.close()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--write-db", action="store_true", help="Also upsert Course/Document/DocumentChunk rows.")
    parser.add_argument("--force", action="store_true", help="Re-ingest even if the manifest says the file is unchanged.")
    args = parser.parse_args()

    args.output.parent.mkdir(parents=True, exist_ok=True)
    manifest_path = args.output.with_suffix(".manifest.json")
    manifest = _load_manifest(manifest_path)
    report = {"files": 0, "chunks": 0, "documents_written": 0, "skipped_unchanged": 0, "errors": []}

    files = sorted(p for p in args.source.rglob("*") if p.suffix.lower() in SUPPORTED)
    with args.output.open("w", encoding="utf-8") as writer:
        for file in files:
            report["files"] += 1
            rel_path = str(file.relative_to(args.source))
            checksum = _file_checksum(file)
            existing = manifest.get(rel_path)
            if not args.force and existing and existing.get("checksum") == checksum and existing.get("status") == "ok":
                report["skipped_unchanged"] += 1
                report["chunks"] += existing.get("chunk_count", 0)
                continue
            try:
                sections = [(location, text.strip()) for location, text in extract(file) if text.strip()]
                for index, (location, text) in enumerate(sections, 1):
                    record = {
                        "id": sha256(f"{file}:{index}".encode()).hexdigest()[:24],
                        "subject_code": subject_code(file),
                        "title": file.stem,
                        "source_path": rel_path,
                        "location": location,
                        "text": text,
                        "checksum": checksum,
                    }
                    writer.write(json.dumps(record, ensure_ascii=False) + "\n")
                    report["chunks"] += 1

                document_id = None
                chunk_count = len(sections)
                if args.write_db and sections:
                    document_id, chunk_count = write_db(
                        rel_path=rel_path,
                        checksum=checksum,
                        subject=subject_code(file),
                        title=file.stem,
                        sections=sections,
                    )
                    report["documents_written"] += 1

                manifest[rel_path] = {
                    "checksum": checksum,
                    "document_id": document_id,
                    "chunk_count": chunk_count,
                    "status": "ok",
                    "error": None,
                }
            except Exception as exc:  # noqa: BLE001 — one bad file must not kill the batch
                report["errors"].append({"path": rel_path, "error": str(exc)})
                manifest[rel_path] = {
                    "checksum": checksum,
                    "document_id": None,
                    "chunk_count": 0,
                    "status": "error",
                    "error": str(exc),
                }

    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    args.output.with_suffix(".report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(
        f"files={report['files']} chunks={report['chunks']} "
        f"documents_written={report['documents_written']} "
        f"skipped_unchanged={report['skipped_unchanged']} errors={len(report['errors'])}"
    )
    if report["errors"]:
        print("Errors:")
        for item in report["errors"]:
            print(f"  {item['path']}: {item['error']}")


if __name__ == "__main__":
    main()
