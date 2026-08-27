"""Ingest pipeline coverage: extraction per format (docx/pptx/xlsx/legacy
.ppt via a mocked LibreOffice), DB write (Document/DocumentChunk), and
manifest-based idempotency (a re-run with unchanged content must not
re-extract or duplicate chunks)."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "scripts"))
import ingest_clean_corpus as ingest  # noqa: E402

from src.db import models  # noqa: E402
from src.db.connection import SessionLocal  # noqa: E402


def _make_docx(path: Path, text: str) -> None:
    doc = DocxDocument()
    doc.add_paragraph(text)
    doc.save(path)


def _make_pptx(path: Path, text: str) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = text
    deck.save(path)


def _make_xlsx(path: Path, rows: list[list[str]]) -> None:
    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(path)


def test_extract_docx_returns_document_section():
    tmp_path = Path("/tmp/test_ingest_docx.docx") if sys.platform != "win32" else Path.cwd() / "_tmp_ingest.docx"
    _make_docx(tmp_path, "Nội dung syllabus CEA201.")
    try:
        sections = ingest.extract(tmp_path)
        assert sections == [("document", "Nội dung syllabus CEA201.")]
    finally:
        tmp_path.unlink(missing_ok=True)


def test_extract_pptx_returns_one_slide_per_entry():
    tmp_path = Path.cwd() / "_tmp_ingest.pptx"
    _make_pptx(tmp_path, "Slide content here")
    try:
        sections = ingest.extract(tmp_path)
        assert len(sections) == 1
        assert sections[0][0] == "slide 1"
        assert "Slide content here" in sections[0][1]
    finally:
        tmp_path.unlink(missing_ok=True)


def test_extract_xlsx_returns_one_sheet_per_entry():
    tmp_path = Path.cwd() / "_tmp_ingest.xlsx"
    _make_xlsx(tmp_path, [["Week", "Topic"], ["1", "Intro"]])
    try:
        sections = ingest.extract(tmp_path)
        assert len(sections) == 1
        assert "Week" in sections[0][1] and "Intro" in sections[0][1]
    finally:
        tmp_path.unlink(missing_ok=True)


def test_extract_ppt_converts_through_mocked_libreoffice(monkeypatch, tmp_path):
    source = tmp_path / "legacy.ppt"
    source.write_bytes(b"not a real ppt, conversion is mocked")
    converted_pptx = tmp_path / "converted.pptx"
    _make_pptx(converted_pptx, "Converted slide text")

    monkeypatch.setattr(ingest.shutil, "which", lambda name: "/usr/bin/soffice")

    def _fake_run(cmd, check, capture_output):
        # Mirror soffice's real behaviour: write "<stem>.pptx" into --outdir.
        outdir = Path(cmd[cmd.index("--outdir") + 1])
        (outdir / f"{source.stem}.pptx").write_bytes(converted_pptx.read_bytes())

    monkeypatch.setattr(ingest.subprocess, "run", _fake_run)

    sections = ingest.extract(source)
    assert sections and "Converted slide text" in sections[0][1]


def test_extract_ppt_without_libreoffice_raises(monkeypatch, tmp_path):
    monkeypatch.setattr(ingest.shutil, "which", lambda name: None)
    with pytest.raises(RuntimeError, match="LibreOffice"):
        ingest.extract(tmp_path / "legacy.ppt")


def test_write_db_creates_course_document_and_chunks():
    sections = [("document", "A" * 2500)]  # forces >1 chunk at CHUNK_SIZE=1100
    document_id, chunk_count = ingest.write_db(
        rel_path="1.ZZTEST/zztest.docx",
        checksum="abc123",
        subject="ZZTEST",
        title="zztest",
        sections=sections,
    )
    assert chunk_count == 3  # 2500 chars / 1100 -> 3 pieces

    db = SessionLocal()
    try:
        course = db.query(models.Course).filter_by(code="ZZTEST").first()
        assert course is not None
        document = db.query(models.Document).filter_by(id=document_id).first()
        assert document is not None
        assert document.metadata_info["source"] == "clean_corpus"
        assert document.metadata_info["original_path"] == "1.ZZTEST/zztest.docx"
        chunks = db.query(models.DocumentChunk).filter_by(document_id=document_id).all()
        assert len(chunks) == 3
    finally:
        db.close()


def test_write_db_is_idempotent_and_replaces_old_chunks():
    document_id, _ = ingest.write_db(
        rel_path="1.ZZTEST2/zztest2.docx", checksum="v1", subject="ZZTEST2", title="zztest2",
        sections=[("document", "first version")],
    )
    document_id_again, chunk_count = ingest.write_db(
        rel_path="1.ZZTEST2/zztest2.docx", checksum="v2", subject="ZZTEST2", title="zztest2",
        sections=[("document", "second version, different text")],
    )
    assert document_id == document_id_again
    assert chunk_count == 1

    db = SessionLocal()
    try:
        chunks = db.query(models.DocumentChunk).filter_by(document_id=document_id).all()
        assert len(chunks) == 1
        assert "second version" in chunks[0].text
    finally:
        db.close()


def test_manifest_skips_unchanged_file_on_rerun(tmp_path):
    source_dir = tmp_path / "clean"
    source_dir.mkdir()
    _make_docx(source_dir / "note.docx", "Same content every time.")
    output = tmp_path / "out" / "corpus.jsonl"

    import subprocess as real_subprocess

    script = str(Path(__file__).resolve().parents[2] / "scripts" / "ingest_clean_corpus.py")
    first = real_subprocess.run(
        [sys.executable, script, "--source", str(source_dir), "--output", str(output)],
        capture_output=True, text=True,
    )
    assert first.returncode == 0, first.stderr
    manifest_path = output.with_suffix(".manifest.json")
    manifest_after_first = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert manifest_after_first["note.docx"]["status"] == "ok"

    second = real_subprocess.run(
        [sys.executable, script, "--source", str(source_dir), "--output", str(output)],
        capture_output=True, text=True,
    )
    assert second.returncode == 0, second.stderr
    assert "skipped_unchanged=1" in second.stdout
